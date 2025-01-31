from pptx import Presentation
from pptx.slide import Slide
from os import listdir, mkdir
from os.path import join
from subprocess import run

# Find pptx file to read from
editor: Presentation
editor_name: str

for file in listdir():
    if file.endswith(".pptx") and not file.startswith("~$"):
        editor = Presentation(file)
        editor_name = file.removesuffix(".pptx")
        break

if editor is None:
    print("No pptx files found")
    exit()


def slide_is_content(slide: Slide) -> bool:
    return slide.shapes[1].name == "Content Placeholder 2"


def text_to_ascii(text: str) -> str:
    """ Replace annoying unicode characters with ascii equivalents """
    ascii_mappings = {
        '“': '"', 
        '”': '"', 
        '‘': "'", 
        '’': "'"
    }
            
    return ''.join(ascii_mappings.get(c, c) for c in text).encode("ascii", "ignore").decode()


project_name: str = editor_name.replace(" ", "_").lower()  # Default project name
title_slide_found: bool = False

# Get project name and description from first slide, if it is a title
if not slide_is_content(editor.slides[0]):
    project_name = editor.slides[0].shapes[0].text.replace(" ", "_").lower() + "_project"
    title_slide_found = True

    # Create project directory
    try:
        mkdir(project_name)
    except Exception: # as e:
        # print(e)
        pass


slide_number = 1
for slide in editor.slides:
    if not slide_is_content(slide):
        # Skip title slides
        continue
    
    # Get file name from title and file content from contents
    slide_title = slide.shapes[0].text
    slide_text = slide.shapes[1].text

    if slide_text == "":
        # Skip empty slides
        continue
    if slide_title == "":
        slide_title = f"{slide_number}.c"

    with open(join(project_name, slide_title), "w") as file:
        file.write(text_to_ascii(slide_text))
        
    slide_number += 1


project_files = [join(project_name, file) for file in listdir(project_name) if file.endswith(".c")]
if len(project_files) > 0:
    # Compile project
    executable = project_name.removesuffix("_project") 
    run(["clang", "-o", executable, *project_files])
    
    # Run executable
    output = run([f"./{executable}"], capture_output=True, text=True)

    if title_slide_found:
        # Put output in title slide
        editor.slides[0].shapes[1].text = output.stdout if output.stdout else output.stderr
        editor.save(editor_name + ".pptx")
